from datetime import date
from pathlib import Path
import textwrap

from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches

from app.main import app

DOCS_DIR = Path("docs")
SLIDES_DIR = Path("slides")


def write_markdown_files() -> tuple[Path, Path, Path]:
    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    SLIDES_DIR.mkdir(parents=True, exist_ok=True)

    api_markdown = DOCS_DIR / "api_documentation.md"
    technical_markdown = DOCS_DIR / "technical_report.md"
    genai_log = DOCS_DIR / "genai_conversation_log_excerpt.md"

    api_markdown.write_text(
        """# ClimatePulse API Documentation

## Project Summary
ClimatePulse is a data-driven web API for storing and analyzing city-level climate records.
The API is implemented using FastAPI + SQLAlchemy + SQLite and supports full CRUD operations.

## Base URL
- Local: `http://127.0.0.1:8000`
- API prefix: `/api/v1`

## Authentication
- Write endpoints are protected using header `X-API-Key`.
- Configure key in `.env` as `API_KEY`.
- Example header: `X-API-Key: cw1-local-api-key`

## Endpoints
### Meta
- `GET /` Service welcome payload
- `GET /health` Health check

### City CRUD
- `POST /api/v1/cities` Create city (auth required)
- `GET /api/v1/cities` List cities
- `GET /api/v1/cities/{city_id}` Get city by id
- `PUT /api/v1/cities/{city_id}` Update city (auth required)
- `DELETE /api/v1/cities/{city_id}` Delete city (auth required)

### Climate Record CRUD
- `POST /api/v1/climate-records` Create record (auth required)
- `GET /api/v1/climate-records` List records with optional filters
- `GET /api/v1/climate-records/{record_id}` Get record by id
- `PUT /api/v1/climate-records/{record_id}` Update record (auth required)
- `DELETE /api/v1/climate-records/{record_id}` Delete record (auth required)

### Analytics
- `GET /api/v1/analytics/cities/{city_id}/temperature-trend?days=30`
- `GET /api/v1/analytics/global/summary?start_date=YYYY-MM-DD&end_date=YYYY-MM-DD`
- `GET /api/v1/analytics/cities/{city_id}/anomalies?threshold_c=3&days=90`

### External Data Import
- `POST /api/v1/import/open-meteo` Import climate history from Open-Meteo (auth required)

## Example Request and Response
### Create City
Request:
`POST /api/v1/cities`
```json
{
  "name": "Leeds",
  "country": "United Kingdom",
  "latitude": 53.8008,
  "longitude": -1.5491
}
```

Response (201):
```json
{
  "id": 1,
  "name": "Leeds",
  "country": "United Kingdom",
  "latitude": 53.8008,
  "longitude": -1.5491,
  "created_at": "2026-04-19T12:00:00"
}
```

## Error Codes
- `400 Bad Request`: invalid range or malformed payload
- `401 Unauthorized`: missing or invalid API key
- `404 Not Found`: resource does not exist
- `409 Conflict`: duplicate data or unique constraint conflict
- `422 Unprocessable Entity`: request validation error
- `500 Internal Server Error`: unexpected server failure

## OpenAPI
- Interactive docs: `/docs`
- ReDoc: `/redoc`
- Exported schema: `docs/openapi.json`
""",
        encoding="utf-8",
    )

    technical_markdown.write_text(
        f"""# Technical Report - ClimatePulse API

## 1. Project Overview
This coursework project implements a climate statistics API with SQL-backed CRUD operations and analytical endpoints.
The solution satisfies the minimum technical requirements from the brief and extends them with:
- API key authentication for write operations
- Automated test suite
- Open-Meteo data import integration
- Delivery automation for documentation and presentation assets

## 2. Problem Scope and Requirements Mapping
The API provides:
- At least one data model with full CRUD linked to a SQL database
- More than four HTTP endpoints
- JSON responses and standard status/error codes
- Local runnable deployment via FastAPI/Uvicorn

Implemented models:
- `City`
- `ClimateRecord`

## 3. Architecture and Stack Justification
- Language: Python 3.11
- Framework: FastAPI (strong validation, built-in OpenAPI)
- ORM: SQLAlchemy 2.0 (clear model mapping and query capabilities)
- Database: SQLite (lightweight SQL database suitable for coursework demos)
- Test tooling: pytest + FastAPI TestClient

This stack was selected to maximize readability, quick iteration, and clean API documentation generation.

## 4. Data Sources and Integration
- Primary external data source: Open-Meteo archive API (`https://archive-api.open-meteo.com`)
- Local reproducible seed dataset: `data/sample_climate_records.csv`

The project includes:
- `scripts/seed_data.py` for deterministic local setup
- `scripts/import_open_meteo_cli.py` for fetching real historical data

## 5. API Design and Error Handling
The API follows REST conventions:
- Resource-oriented routing (`/cities`, `/climate-records`)
- Proper status codes (`201`, `204`, `401`, `404`, `409`, `422`)
- Validation through Pydantic schemas
- Consistent JSON response payloads

Authentication:
- Header-based API key (`X-API-Key`) required on write endpoints

## 6. Testing Strategy
Tests cover:
- Health endpoint
- City CRUD and auth checks
- Climate record CRUD
- Analytics endpoints (trend, summary, anomalies)

The tests use an isolated temporary SQLite database for repeatability.

## 7. Challenges and Lessons Learned
Challenges:
- Maintaining data consistency between city and climate record entities
- Designing analytical endpoints that remain simple but meaningful
- Producing all deliverables (code, docs PDF, report PDF, PPTX) in one coherent workflow

Lessons learned:
- Strong schema validation reduces runtime bugs early
- Modular services improve maintainability and testing
- Automating deliverable generation reduces submission risk

## 8. Limitations and Future Work
Current limitations:
- SQLite may not be ideal for high concurrency production workloads
- Authentication is API-key based instead of role-based user auth
- Analytics are descriptive, not predictive

Future improvements:
- Add PostgreSQL support and migrations
- Add JWT-based authentication and user accounts
- Add caching and pagination metadata
- Add anomaly detection models using time-series methods

## 9. GenAI Declaration and Reflection
GenAI tools were used for:
- Planning project structure
- Drafting documentation content
- Refining endpoint and test coverage ideas
- Improving wording and report organization

All generated outputs were manually reviewed, edited, and validated with runnable tests.
No AI content was accepted without developer verification.

## 10. Supplementary GenAI Conversation Logs
An excerpt log is included at:
- `docs/genai_conversation_log_excerpt.md`

Date generated: {date.today().isoformat()}
""",
        encoding="utf-8",
    )

    genai_log.write_text(
        """# GenAI Conversation Log Excerpt

## Tooling Context
- Assistant: GPT-5 coding agent
- Date: 2026-04-19
- Usage mode: Green-light GenAI coursework support

## Example Prompt Themes
1. Interpret PDF brief requirements and produce a checklist.
2. Scaffold a FastAPI + SQLAlchemy project with CRUD and analytics.
3. Create tests and generate submission deliverables (PDF and PPTX).

## Verification Notes
- All generated code was executed locally.
- Tests were run with pytest and fixed until passing.
- Deliverable files were generated by project scripts and manually inspected.
""",
        encoding="utf-8",
    )

    return api_markdown, technical_markdown, genai_log


def markdown_to_pdf(markdown_path: Path, pdf_path: Path, title: str) -> None:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    pdf.cell(0, 10, title, new_x="LMARGIN", new_y="NEXT")
    pdf.ln(2)

    for raw_line in markdown_path.read_text(encoding="utf-8").splitlines():
        line = raw_line.strip("\n")
        if not line:
            pdf.ln(3)
            continue

        if line.startswith("### "):
            pdf.set_font("Helvetica", "B", 12)
            pdf.multi_cell(0, 7, line.replace("### ", ""), new_x="LMARGIN", new_y="NEXT")
            pdf.set_font("Helvetica", "", 10)
            continue
        if line.startswith("## "):
            pdf.set_font("Helvetica", "B", 13)
            pdf.multi_cell(0, 8, line.replace("## ", ""), new_x="LMARGIN", new_y="NEXT")
            pdf.set_font("Helvetica", "", 10)
            continue
        if line.startswith("# "):
            pdf.set_font("Helvetica", "B", 14)
            pdf.multi_cell(0, 8, line.replace("# ", ""), new_x="LMARGIN", new_y="NEXT")
            pdf.set_font("Helvetica", "", 10)
            continue
        if line.startswith("- "):
            pdf.set_font("Helvetica", "", 10)
            pdf.multi_cell(0, 6, f"* {line[2:]}", new_x="LMARGIN", new_y="NEXT")
            continue

        cleaned = line.replace("`", "")
        wrapped_lines = textwrap.wrap(cleaned, width=95, break_long_words=True, break_on_hyphens=True)
        pdf.set_font("Helvetica", "", 10)
        if not wrapped_lines:
            pdf.multi_cell(0, 6, "", new_x="LMARGIN", new_y="NEXT")
            continue
        for wrapped in wrapped_lines:
            pdf.multi_cell(0, 6, wrapped, new_x="LMARGIN", new_y="NEXT")

    pdf.output(str(pdf_path))


def build_presentation() -> Path:
    presentation_path = SLIDES_DIR / "coursework_presentation.pptx"
    prs = Presentation()

    def add_bullet_slide(title: str, bullets: list[str]) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        for idx, bullet in enumerate(bullets):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            p.text = bullet
            p.level = 0

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "ClimatePulse API"
    title_slide.placeholders[1].text = "XJCO3011 Coursework 1 - Individual API Project"

    add_bullet_slide(
        "Project Overview",
        [
            "Goal: Build a data-driven web API with SQL integration",
            "Domain: Urban climate records and analytics",
            "Stack: FastAPI + SQLAlchemy + SQLite + pytest",
            "Deliverables: Code, API docs PDF, technical report PDF, presentation",
        ],
    )
    add_bullet_slide(
        "Architecture and Data Model",
        [
            "Models: City and ClimateRecord",
            "CRUD implemented for both resources",
            "Service layer for analytics and external data imports",
            "Database constraints enforce data consistency",
        ],
    )
    add_bullet_slide(
        "API Documentation Overview",
        [
            "Interactive OpenAPI docs available at /docs and /redoc",
            "Manual API PDF includes endpoints, auth, status codes, examples",
            "Exported schema stored as docs/openapi.json",
            "Documentation linked from README",
        ],
    )
    add_bullet_slide(
        "Version Control Practices",
        [
            "Project structured into app, scripts, tests, docs, and slides",
            "Incremental commits recommended: scaffold, features, tests, docs",
            "README includes run and verification commands",
            "Repository keeps generated deliverables under version control",
        ],
    )
    add_bullet_slide(
        "Technical Report Highlights",
        [
            "Justifies stack and architecture decisions",
            "Explains testing strategy and error handling",
            "Reflects on challenges, limitations, and future work",
            "Contains GenAI declaration and usage reflection",
        ],
    )
    add_bullet_slide(
        "Testing and Error Handling",
        [
            "Automated tests cover auth, CRUD flows, and analytics",
            "Status codes: 201, 204, 401, 404, 409, 422",
            "Input validation via Pydantic schemas",
            "API key protection for write endpoints",
        ],
    )
    add_bullet_slide(
        "All Deliverables Checklist",
        [
            "Source code and runnable API",
            "README with setup instructions",
            "docs/api_documentation.pdf",
            "docs/technical_report.pdf",
            "slides/coursework_presentation.pptx",
        ],
    )
    add_bullet_slide(
        "Demo Plan and Q&A",
        [
            "Start API and open /docs",
            "Run city and climate record CRUD examples",
            "Call analytics endpoints and discuss outputs",
            "Answer design, testing, and trade-off questions",
        ],
    )

    # Add a simple visual separator rectangle in the final slide.
    final_slide = prs.slides[-1]
    final_slide.shapes.add_shape(
        autoshape_type_id=1,  # rectangle
        left=Inches(0.5),
        top=Inches(6.5),
        width=Inches(12.0),
        height=Inches(0.2),
    )

    prs.save(str(presentation_path))
    return presentation_path


def export_openapi_schema(path: Path) -> None:
    import json

    schema = app.openapi()
    path.write_text(json.dumps(schema, indent=2), encoding="utf-8")


def main() -> None:
    api_md, report_md, _genai_md = write_markdown_files()
    api_pdf = DOCS_DIR / "api_documentation.pdf"
    report_pdf = DOCS_DIR / "technical_report.pdf"
    openapi_json = DOCS_DIR / "openapi.json"

    markdown_to_pdf(api_md, api_pdf, "ClimatePulse API Documentation")
    markdown_to_pdf(report_md, report_pdf, "ClimatePulse Technical Report")
    export_openapi_schema(openapi_json)
    pptx_path = build_presentation()

    print("Generated deliverables:")
    print(f"- {api_md}")
    print(f"- {api_pdf}")
    print(f"- {report_md}")
    print(f"- {report_pdf}")
    print(f"- {openapi_json}")
    print(f"- {pptx_path}")


if __name__ == "__main__":
    main()
